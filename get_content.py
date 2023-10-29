# get_content.py

import sys
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import os
import pandas as pd
import string

class FileContentReader:
    def __init__(self):
        self.file_paths = []

    def set_file_paths(self, file_paths):
        self.file_paths = file_paths

    def read_pdf(self, file_path):
        text = ""
        pdf_reader = PdfReader(file_path)
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text

    def read_docx(self, file_path):
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
        return text

    def read_pptx(self, file_path):
        ppt = Presentation(file_path)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text

    def read_txt(self, file_path):
        with open(file_path, 'r') as txt_file:
            text = txt_file.read()
        return text

    def read_content(self):
        combined_content = ""
        file_extension=""
        for file_path in self.file_paths:
            file_extension = os.path.splitext(file_path)[1].lower()
            if file_extension == ".pdf":
                combined_content += self.read_pdf(file_path)
            elif file_extension == ".docx":
                combined_content += self.read_docx(file_path)
            elif file_extension == ".pptx":
                combined_content += self.read_pptx(file_path)
            elif file_extension == ".txt":
                combined_content += self.read_txt(file_path)
            elif file_extension == ".csv":
                df = pd.read_csv(file_path)
                return df.astype(str), "google/tapas-base-finetuned-wtq"

        return combined_content, "deepset/bert-base-cased-squad2"

    def get_model_name(self, file_extension):
        if file_extension == ".csv":
            return "google/tapas-base-finetuned-wtq"
        else:
            return "deepset/bert-base-cased-squad2"

    def process_files(self):
        combined_content, modelname = self.read_content()
        return combined_content, modelname

